import os
from pptx import Presentation
from pptx.util import Inches, Pt, Cm
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from PIL import Image, ImageDraw

# --- CONFIGURACIÓN DE RUTAS Y NOMBRES ---
OUTPUT_FILENAME = "Presentacion_Convocatoria_46_Estructurada.pptx"
IMG_DIR = "assets_temp"
LOGO_PATH = os.path.join("CTM_Agents/static", "CotecmarLogo.png")
COVER_PATH = os.path.join("CTM_Agents/generated_images", "detección_de_anomalías_sísmicas_con_inteligencia_artificial.png")


# --- PALETA DE COLORES ---
COLOR_PRIMARY = RGBColor(0, 51, 102)     # Azul Oscuro Corporativo
COLOR_ACCENT = RGBColor(255, 192, 0)     # Amarillo Acento
COLOR_WHITE = RGBColor(255, 255, 255)
COLOR_TEXT_DARK = RGBColor(42, 42, 42)
COLOR_TABLE_HEADER = RGBColor(0, 76, 153)

def create_dummy_assets():
    """Genera imágenes placeholder si no existen."""
    if not os.path.exists(IMG_DIR):
        os.makedirs(IMG_DIR)
    
    if not os.path.exists(LOGO_PATH):
        img = Image.new('RGB', (200, 80), color=(255, 255, 255))
        d = ImageDraw.Draw(img)
        d.text((10, 30), "LOGO ENTIDAD", fill=(0, 51, 102))
        img.save(LOGO_PATH)

    if not os.path.exists(COVER_PATH):
        img = Image.new('RGB', (1280, 720), color=(0, 51, 102))
        d = ImageDraw.Draw(img)
        d.text((500, 300), "IMAGEN DE FONDO", fill=(255, 255, 255))
        img.save(COVER_PATH)

def add_header(slide, title_text):
    """Agrega el encabezado estético (Azul + Amarillo + Logo)."""
    width = Presentation().slide_width
    header_height = Cm(2.8)
    
    # 1. Acento Amarillo (Izquierda)
    yellow_width = Cm(1.2)
    shape_yellow = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, yellow_width, header_height)
    shape_yellow.fill.solid()
    shape_yellow.fill.fore_color.rgb = COLOR_ACCENT
    shape_yellow.line.fill.background()

    # 2. Barra Azul (Principal)
    shape_blue = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, yellow_width, 0, width - yellow_width, header_height)
    shape_blue.fill.solid()
    shape_blue.fill.fore_color.rgb = COLOR_PRIMARY
    shape_blue.line.fill.background()

    # 3. Título
    text_box = slide.shapes.add_textbox(yellow_width + Cm(0.5), Cm(0), width - yellow_width - Cm(4), header_height)
    tf = text_box.text_frame
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.text = title_text.upper()
    p.font.name = 'Arial'
    p.font.size = Pt(22)
    p.font.bold = True
    p.font.color.rgb = COLOR_WHITE

    # 4. Logo
    if os.path.exists(LOGO_PATH):
        slide.shapes.add_picture(LOGO_PATH, width - Cm(4.5), Cm(0.4), width=Cm(4))

def add_bullet_text(slide, lines, font_size=16, bold_first=False):
    """Agrega una caja de texto con bullets."""
    left = Cm(1.5)
    top = Cm(3.5)
    width = Presentation().slide_width - Cm(3)
    height = Presentation().slide_height - Cm(4)
    
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    
    for line in lines:
        p = tf.add_paragraph()
        
        # Lógica de indentación simple basada en caracteres
        clean_text = line
        level = 0
        if line.startswith("•") or line.startswith("-"):
            level = 1
            clean_text = line.lstrip("•- ")
        elif line.startswith("  "):
            level = 2
            clean_text = line.strip()

        p.text = clean_text
        p.level = level
        p.font.size = Pt(font_size)
        p.font.color.rgb = COLOR_TEXT_DARK
        p.space_after = Pt(6)
        
        if bold_first and level == 0:
            p.font.bold = True

def add_table(slide, data, col_widths):
    """Crea una tabla estilizada."""
    rows = len(data)
    cols = len(data[0])
    left = Cm(1.5)
    top = Cm(4.0)
    width = Presentation().slide_width - Cm(3)
    
    table_shape = slide.shapes.add_table(rows, cols, left, top, width, Cm(0.8 * rows)).table
    
    # Ajustar anchos
    for i, w in enumerate(col_widths):
        table_shape.columns[i].width = Cm(w)

    for r in range(rows):
        for c in range(cols):
            cell = table_shape.cell(r, c)
            cell.text = str(data[r][c])
            p = cell.text_frame.paragraphs[0]
            p.font.size = Pt(12)
            
            if r == 0: # Header
                cell.fill.solid()
                cell.fill.fore_color.rgb = COLOR_TABLE_HEADER
                p.font.color.rgb = COLOR_WHITE
                p.font.bold = True
                p.alignment = PP_ALIGN.CENTER
            else: # Body
                p.font.color.rgb = COLOR_TEXT_DARK
                cell.fill.solid()
                cell.fill.fore_color.rgb = RGBColor(245, 245, 245) if r % 2 == 0 else COLOR_WHITE

def create_presentation():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    
    # --- 1. PORTADA ---
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    if os.path.exists(COVER_PATH):
        slide.shapes.add_picture(COVER_PATH, 0, 0, width=prs.slide_width, height=prs.slide_height)
    
    # Banner inferior portada
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, prs.slide_height - Cm(5), prs.slide_width, Cm(5))
    shape.fill.solid()
    shape.fill.fore_color.rgb = COLOR_PRIMARY
    shape.fill.alpha = 0.9
    shape.line.fill.background()
    
    tb = slide.shapes.add_textbox(Cm(1), prs.slide_height - Cm(4.5), prs.slide_width - Cm(2), Cm(4))
    p = tb.text_frame.paragraphs[0]
    p.text = "CONVOCATORIA 46 COLOMBIA INTELIGENTE\nINFRAESTRUCTURA PARA EL DESARROLLO DE LA INTELIGENCIA ARTIFICIAL"
    p.font.bold = True
    p.font.size = Pt(28)
    p.font.color.rgb = COLOR_WHITE
    p.alignment = PP_ALIGN.CENTER

    # --- 1. DATOS GENERALES ---
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_header(slide, "1. Datos Generales y Fechas")
    lines = [
        "Fechas Relevantes:",
        "• Apertura: Martes 14 de octubre de 2025",
        "• Cierre: Viernes 14 de noviembre de 2025",
        "• Resultados Preliminares: Lunes 29 de diciembre de 2025",
        "• Resultados Definitivos: Lunes 29 de diciembre de 2025",
        "",
        "Nota: Los plazos son improrrogables según términos de referencia."
    ]
    add_bullet_text(slide, lines, font_size=20, bold_first=True)

    # --- 2. OBJETIVO ---
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_header(slide, "2. Objetivo General")
    lines = [
        "Propósito:",
        "Impulsar la infraestructura, el desarrollo científico y tecnológico en inteligencia artificial.",
        "",
        "Elementos Clave:",
        "• Creación y consolidación de capacidades nacionales avanzadas.",
        "• Habilitar la investigación, el desarrollo y la innovación.",
        "• Fortalecer la soberanía tecnológica del país.",
        "• Posicionar al país como referente regional en IA."
    ]
    add_bullet_text(slide, lines, font_size=20, bold_first=True)

    # --- 3. DIRIGIDO A (Parte 1) ---
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_header(slide, "3. Dirigido A (Alianza Estratégica)")
    lines = [
        "Composición Mínima de la Alianza:",
        "• Al menos UN (1) actor del SNCTI reconocido.",
        "• Al menos DOS (2) entidades territoriales (1 Gobernación + 1 Alcaldía).",
        "• Al menos UNA (1) Empresa Nacional.",
        "",
        "Requisitos de Desempeño:",
        "• Grupo de Investigación A1 o A (con líneas en IA).",
        "• El proponente debe acreditar adecuado desempeño según medición DNP.",
        "• Valor mínimo de contrapartida: 10% de recursos solicitados."
    ]
    add_bullet_text(slide, lines, font_size=18, bold_first=True)

    # --- 3. DIRIGIDO A (Parte 2 - Experiencia) ---
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_header(slide, "3. Experiencia Requerida (Últimos 10 años)")
    lines = [
        "Total Proyectos Requeridos: 5 Proyectos CTeI ejecutados/en ejecución",
        "",
        "Distribución Temática:",
        "• Mínimo 3 proyectos en:",
        "  - Inteligencia Artificial / Modelos de Lenguaje",
        "  - Ciencia de Datos / Supercómputo",
        "• Mínimo 2 proyectos en:",
        "  - Implementación, construcción o desarrollo de infraestructura científica/tecnológica.",
        "",
        "Notas Adicionales:",
        "• Se permite participación de entidades internacionales.",
        "• Roles definidos: Ejecutor vs Aliados.",
        "• Restricción: Máximo 1 proyecto presentado por entidad."
    ]
    add_bullet_text(slide, lines, font_size=18, bold_first=True)

    # --- 4. DEMANDAS TERRITORIALES ---
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_header(slide, "4. Demandas Territoriales")
    lines = [
        "Enfoque Regional:",
        "• Las propuestas deben atender necesidades regionales específicas.",
        "• Consultar Anexo 4: Demandas Territoriales Bienio 2025-2026.",
        "• Se requiere articulación con los CODECTI departamentales.",
        "",
        "Requisito de Cobertura:",
        "• El proyecto debe cubrir al menos una demanda por cada región del SGR."
    ]
    add_bullet_text(slide, lines, font_size=20, bold_first=True)

    # --- 5. LÍNEAS TEMÁTICAS (Tabla) ---
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_header(slide, "5. Líneas Temáticas")
    
    data_lineas = [
        ["#", "Línea Temática", "Enfoque Principal"],
        ["1", "Infraestructura tecnológica y científica", "Hardware, procesamiento, almacenamiento y redes."],
        ["2", "Desarrollo de modelos avanzados de IA", "LLMs, Modelos fundacionales para sectores estratégicos."],
        ["3", "Formación de talento humano", "Capacitación especializada, atracción de expertos (PhD)."],
        ["4", "Propiedad intelectual y transferencia", "Patentes, spin-offs, licenciamiento de tecnología."],
        ["5", "Articulación entre actores", "Conexión efectiva SNCTI + Sector Productivo + Estado."]
    ]
    add_table(slide, data_lineas, [2, 12, 14])

    # --- 6. ALCANCE DEL PROYECTO ---
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_header(slide, "6. Alcance del Proyecto")
    lines = [
        "Componentes Obligatorios:",
        "✅ Infraestructura tecnológica avanzada (HPC, Data Centers).",
        "✅ Desarrollo de modelos de IA de gran escala.",
        "✅ Fortalecimiento del talento humano (Becas, Estancias).",
        "✅ Innovación y transferencia tecnológica.",
        "✅ Articulación interinstitucional.",
        "",
        "Cobertura de Demandas:",
        "• Mínimo 6 demandas territoriales atendidas.",
        "• Distribución: Al menos 1 demanda por cada región SGR."
    ]
    add_bullet_text(slide, lines, font_size=18, bold_first=True)

    # --- 7. PRODUCTOS E INDICADORES ---
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_header(slide, "7. Productos e Indicadores de CTeI")
    lines = [
        "Generación de Nuevo Conocimiento:",
        "• Cód 3906012: Documentos de investigación (Artículos A1, A2).",
        "• Cód 3906013: Prototipos industriales/tecnológicos.",
        "",
        "Uso y Transferencia:",
        "• Cód 3906009: Apoyo financiero a programas CTI.",
        "• Cód 3906014: Servicios de asistencia técnica.",
        "",
        "Apropiación Social:",
        "• Cód 3906016: Documentos de política pública.",
        "• Cód 3906015: Documentos de planeación estratégica.",
        "",
        "Infraestructura:",
        "• Cód 3906018/020/021: Infraestructura construida, dotada y adecuada."
    ]
    add_bullet_text(slide, lines, font_size=16, bold_first=True)

    # --- 8. CONSIDERACIONES TÉCNICAS ---
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_header(slide, "8. Consideraciones Técnicas")
    lines = [
        "Aspectos Técnicos:",
        "• Escalabilidad y diseño modular.",
        "• Soporte específico a cargas de trabajo de IA.",
        "• Confiabilidad y durabilidad mínima de 10 años.",
        "",
        "Estándares y Normatividad:",
        "• Cumplimiento del estándar Tier III (Disponibilidad).",
        "• Eficiencia energética y sostenibilidad ambiental.",
        "• Ciberseguridad y alta disponibilidad.",
        "",
        "Requisito General:",
        "• Alcance nacional e integralidad de la solución."
    ]
    add_bullet_text(slide, lines, font_size=18, bold_first=True)

    # --- 9. ENFOQUES ---
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_header(slide, "9. Enfoques Territorial y Diferencial")
    lines = [
        "Enfoque Territorial (Puntaje Adicional):",
        "• Priorización de proyectos ejecutados fuera de los grandes centros:",
        "  - Bogotá / Cundinamarca",
        "  - Antioquia",
        "  - Valle del Cauca",
        "  - Atlántico",
        "",
        "Enfoque Diferencial (Inclusión Social):",
        "• Inclusión explícita de:",
        "  - Población víctima del conflicto armado.",
        "  - Personas con discapacidad.",
        "  - Grupos étnicos.",
        "  - Mujeres (Paridad de género)."
    ]
    add_bullet_text(slide, lines, font_size=19, bold_first=True)

    # --- 10. TALENTO HUMANO ---
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_header(slide, "10. Vinculación de Talento Humano")
    lines = [
        "Equipo Mínimo Obligatorio a Financiar:",
        "",
        "1. Nivel Postdoctoral:",
        "• Mínimo 1 Estancia posdoctoral.",
        "",
        "2. Nivel Maestría:",
        "• Mínimo 2 Estudiantes de maestría vinculados.",
        "",
        "3. Jóvenes Investigadores:",
        "• Total: 6 Jóvenes.",
        "• Distribución: 3 de pregrado + 3 profesionales."
    ]
    add_bullet_text(slide, lines, font_size=20, bold_first=True)

    # --- 11. REQUISITOS ---
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_header(slide, "11. Requisitos Documentales (Resumen)")
    lines = [
        "1. Inscripción en plataforma SIGP.",
        "2. Verificación de requisitos 'Dirigido a'.",
        "3. Certificado de adecuado desempeño (DNP).",
        "4. Carta de Aval y Modelo de Gobernanza.",
        "5. Documento técnico completo.",
        "6. Presupuesto detallado.",
        "7. Proyecto cargado en MGA (borrador).",
        "8. Certificados de contrapartida.",
        "9. Cartas de participación de aliados.",
        "10. Soportes de experiencia (5 proyectos).",
        "11. Aval Comité de Ética (si aplica).",
        "12. Acreditación de Alta Calidad IES.",
        "13. Resolución de reconocimiento de Centros.",
        "14-17. Documentos legales y atención a demandas."
    ]
    add_bullet_text(slide, lines, font_size=14, bold_first=False)

    # --- 12. DURACIÓN Y FINANCIACIÓN ---
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_header(slide, "12. Duración y Financiación")
    lines = [
        "Duración del Proyecto:",
        "• Plazo de ejecución: Hasta 48 meses.",
        "",
        "Recursos Disponibles:",
        "• Fuente: Sistema General de Regalías (Asignación CTI).",
        "• Monto Total de la Convocatoria: $630.000.000.000 COP.",
        "",
        "Expectativa de Adjudicación:",
        "• Se espera financiar UN (1) único proyecto de alcance nacional."
    ]
    add_bullet_text(slide, lines, font_size=20, bold_first=True)

    # --- 13. RIESGOS (Críticos - Tabla) ---
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_header(slide, "13. Mapa de Riesgos: Críticos 🔴")
    
    # Tabla de riesgos críticos extraída del texto
    data_riesgos = [
        ["#", "Riesgo", "Categoría", "Estrategia Mitigación"],
        ["1", "Inexperiencia Proyectos SGR", "Institucional", "Apoyo UNAL, talleres constantes"],
        ["2", "Inexperiencia GESPROY", "Administrativo", "Capacitación, equipo mixto"],
        ["3", "Capacidad Admin. Insuficiente", "Financiero", "Unidad mixta de gestión"],
        ["5", "Restricción Nómina", "Financiero", "Contratación por producto/servicio"],
        ["7", "Personal no certificado", "Gobernanza", "Personal certificado SGR"],
        ["9", "Dependencia Proveedores Int.", "Externo", "Diversificar proveedores"],
        ["13", "Brechas Ciberseguridad", "Técnico", "ISO 27001, auditorías"],
        ["16", "Débil Gobernanza", "Estratégico", "Modelo tripartito formal"]
    ]
    # Ajuste para que quepa en la diapositiva
    add_table(slide, data_riesgos, [1.5, 9, 5, 12])

    # --- 13. RIESGOS (Altos y Moderados) ---
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_header(slide, "13. Mapa de Riesgos: Altos y Moderados")
    lines = [
        "Riesgos Altos 🟡 (8 riesgos):",
        "• Financiación supervisor/interventor.",
        "• Sobrecostos en implementación tecnológica.",
        "• Falta de articulación entre aliados.",
        "• Fallas en sostenibilidad post-ejecución.",
        "• Riesgos reputacionales y de control fiscal.",
        "",
        "Riesgos Moderados 🟢 (1 riesgo):",
        "• Falta de alineación con estándares internacionales de IA ética.",
        "",
        "Distribución Porcentual:",
        "• Críticos: 47% | Altos: 47% | Moderados: 6%"
    ]
    add_bullet_text(slide, lines, font_size=18, bold_first=True)

    # --- GUARDAR ---
    prs.save(OUTPUT_FILENAME)
    print(f"✅ Presentación generada exitosamente: {OUTPUT_FILENAME}")

if __name__ == "__main__":
    create_dummy_assets()
    create_presentation()